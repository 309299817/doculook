# Copyright (c) Opendatalab. All rights reserved.
"""
Office文档转换器模块
支持将Word、Excel、PPT文档转换为PDF格式
"""

import io
import tempfile
import os
import sys
import subprocess
import shutil
from platform import system as platform_system
from pathlib import Path
from typing import Union, BinaryIO
from loguru import logger

try:
    import docx
    from docx import Document
    from docx.document import Document as DocumentType
except ImportError:
    logger.warning("python-docx 未安装，Word文档支持将被禁用")
    docx = None
    Document = None
    DocumentType = None

try:
    import openpyxl
    from openpyxl import Workbook
    import pandas as pd
except ImportError:
    logger.warning("openpyxl 或 pandas 未安装，Excel文档支持将被禁用")
    openpyxl = None
    Workbook = None
    pd = None

try:
    from pptx import Presentation
except ImportError:
    logger.warning("python-pptx 未安装，PowerPoint文档支持将被禁用")
    Presentation = None

try:
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
except ImportError:
    logger.error("reportlab 未安装，无法进行PDF转换")
    raise ImportError("缺少必需的依赖: reportlab")

try:
    from PIL import Image, ImageEnhance
except ImportError:
    logger.warning("PIL (Pillow) 未安装，图片处理功能将受限")
    Image = None
    ImageEnhance = None


def _extract_image_from_drawing(drawing, doc):
    """
    从Word文档的drawing元素中提取图片数据
    
    Args:
        drawing: Word文档中的drawing元素
        doc: Word文档对象
        
    Returns:
        bytes: 图片数据，如果提取失败返回None
    """
    try:
        # 查找图片关系ID
        blip_elements = drawing.xpath('.//a:blip')
        if not blip_elements:
            return None
        
        blip = blip_elements[0]
        embed_id = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
        
        if embed_id:
            # 从文档关系中获取图片数据
            image_part = doc.part.related_parts[embed_id]
            return image_part.blob
        
    except Exception as e:
        logger.debug(f"提取图片数据失败: {e}")
    
    return None


def _extract_images_from_word_document(doc):
    """
    从Word文档中提取所有图片
    
    Args:
        doc: Word文档对象
        
    Returns:
        list: 图片数据列表
    """
    images = []
    try:
        # 方法1: 从文档关系中提取图片
        for rel in doc.part.rels.values():
            if hasattr(rel, 'target_ref') and "image" in rel.target_ref:
                try:
                    image_data = rel.target_part.blob
                    images.append(image_data)
                except Exception as e:
                    logger.debug(f"提取关系图片失败: {e}")
        
        # 方法2: 从document.xml中查找图片引用
        if hasattr(doc, '_element'):
            try:
                # 查找所有的图片引用
                for blip in doc._element.xpath('.//a:blip[@r:embed]'):
                    embed_id = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    if embed_id and embed_id in doc.part.rels:
                        try:
                            image_part = doc.part.rels[embed_id].target_part
                            image_data = image_part.blob
                            if image_data not in images:  # 避免重复
                                images.append(image_data)
                        except Exception as e:
                            logger.debug(f"提取blip图片失败: {e}")
            except Exception as e:
                logger.debug(f"解析document.xml失败: {e}")
                
    except Exception as e:
        logger.debug(f"提取Word图片失败: {e}")
    
    logger.info(f"从Word文档中提取到 {len(images)} 张图片")
    return images


def _enhance_image_for_ocr(pil_image):
    """
    增强图片质量以提高OCR识别效果
    
    Args:
        pil_image: PIL图片对象
        
    Returns:
        PIL.Image: 增强后的图片
    """
    if not ImageEnhance:
        return pil_image
    
    try:
        # 转换为RGB模式以确保兼容性
        if pil_image.mode not in ('RGB', 'L'):
            pil_image = pil_image.convert('RGB')
        
        # 增强对比度以提高文字清晰度
        enhancer = ImageEnhance.Contrast(pil_image)
        pil_image = enhancer.enhance(1.2)  # 增强对比度20%
        
        # 增强锐度以提高边缘清晰度
        enhancer = ImageEnhance.Sharpness(pil_image)
        pil_image = enhancer.enhance(1.1)  # 增强锐度10%
        
        return pil_image
    except Exception as e:
        logger.debug(f"图片增强失败: {e}")
        return pil_image


def extract_word_content_and_images(word_bytes: bytes) -> dict:
    """
    直接从Word文档中提取文本内容和图片
    
    Args:
        word_bytes: Word文档的字节数据
        
    Returns:
        dict: 包含文本内容和图片数据的字典
    """
    if Document is None:
        raise ImportError("python-docx 未安装，无法处理Word文档")
    
    try:
        # 使用BytesIO读取Word文档
        word_stream = io.BytesIO(word_bytes)
        doc = Document(word_stream)
        
        # 提取文本内容
        text_content = []
        
        # 按顺序处理文档中的所有元素
        table_index = 0
        image_index = 0
        
        # 遍历文档body中的所有元素
        for element in doc.element.body:
            # 处理段落元素
            if element.tag.endswith('}p'):  # 段落元素
                # 查找对应的段落对象
                for paragraph in doc.paragraphs:
                    if paragraph._element == element:
                        if paragraph.text.strip():
                            # 根据段落样式添加标记
                            if paragraph.style.name.startswith('Heading'):
                                text_content.append(f"# {paragraph.text}")
                            else:
                                text_content.append(paragraph.text)
                        
                        # 检查段落中是否包含图片
                        for run in paragraph.runs:
                            for drawing in run._element.xpath('.//w:drawing'):
                                image_index += 1
                                # 在图片位置插入图片引用
                                image_ref = f"[图片{image_index}]"
                                text_content.append(image_ref)
                        break
            
            # 处理表格元素
            elif element.tag.endswith('}tbl'):  # 表格元素
                # 查找对应的表格对象
                for table in doc.tables:
                    if table._element == element:
                        table_index += 1
                        table_text = []
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                table_text.append(" | ".join(row_text))
                        
                        if table_text:
                            text_content.append(f"\n**表格{table_index}:**")
                            text_content.extend(table_text)
                            text_content.append("")
                        break
        
        # 提取图片数据
        document_images = _extract_images_from_word_document(doc)
        
        return {
            "text_content": "\n\n".join(text_content),
            "images": document_images,
            "has_text": len([t for t in text_content if t.strip()]) > 0,
            "has_images": len(document_images) > 0
        }
            
    except Exception as e:
        logger.error(f"Word内容提取失败: {e}")
        raise


def word_to_pdf_bytes(word_bytes: bytes) -> bytes:
    """
    处理Word文档：如果有文本内容就直接使用，只对图片转换为PDF供OCR使用
    
    Args:
        word_bytes: Word文档的字节数据
        
    Returns:
        PDF文档的字节数据（仅包含需要OCR的图片）
    """
    try:
        # 提取Word内容
        word_content = extract_word_content_and_images(word_bytes)
        
        # 如果有文本内容，记录到日志中供后续处理使用
        if word_content["has_text"]:
            logger.info(f"Word文档包含文本内容，长度: {len(word_content['text_content'])}")
            # 将文本内容存储到全局变量或其他方式，供后续处理使用
            # 这里暂时记录到日志，实际应该有更好的传递方式
        
        # 如果没有图片，返回空PDF或者抛出异常表示不需要OCR
        if not word_content["has_images"]:
            logger.info("Word文档不包含图片，不需要创建PDF进行OCR")
            # 返回一个最小的空PDF
            pdf_buffer = io.BytesIO()
            pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=A4)
            story = [Paragraph("No images to process", getSampleStyleSheet()['Normal'])]
            pdf_doc.build(story)
            return pdf_buffer.getvalue()
        
        # 只为图片创建PDF供OCR使用
        pdf_buffer = io.BytesIO()
        pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=A4)
        story = []
        
        logger.info(f"处理Word文档中的 {len(word_content['images'])} 张图片")
        
        for i, image_data in enumerate(word_content['images']):
            try:
                if Image:
                    image_stream = io.BytesIO(image_data)
                    pil_image = Image.open(image_stream)
                    
                    # 保持较高的图片分辨率以确保OCR质量
                    max_width = 1200
                    max_height = 900
                    
                    # 只在图片过大时才缩放
                    if pil_image.width > max_width or pil_image.height > max_height:
                        pil_image.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)
                    
                    # 增强图片质量以提高OCR效果
                    pil_image = _enhance_image_for_ocr(pil_image)
                    
                    # 保存增强后的图片
                    resized_stream = io.BytesIO()
                    pil_image.save(resized_stream, format='PNG', quality=95, optimize=True)
                    resized_stream.seek(0)
                    
                    # 计算适合PDF页面的图片尺寸
                    page_width = 6*inch
                    page_height = 8*inch
                    
                    width_ratio = page_width / pil_image.width
                    height_ratio = page_height / pil_image.height
                    scale_ratio = min(width_ratio, height_ratio, 1.0)
                    
                    final_width = pil_image.width * scale_ratio
                    final_height = pil_image.height * scale_ratio
                    
                    # 添加到PDF
                    rl_image = RLImage(resized_stream, width=final_width, height=final_height)
                    story.append(rl_image)
                    story.append(Spacer(1, 0.2*inch))
                    
            except Exception as e:
                logger.warning(f"处理图片 {i+1} 失败: {e}")
        
        # 如果没有成功处理的图片，添加一个占位符
        if not story:
            story.append(Paragraph("No processable images found", getSampleStyleSheet()['Normal']))
        
        # 生成PDF
        pdf_doc.build(story)
        return pdf_buffer.getvalue()
            
    except Exception as e:
        logger.error(f"Word处理失败: {e}")
        raise


def excel_to_pdf_bytes(excel_bytes: bytes) -> bytes:
    """
    将Excel文档字节转换为PDF字节
    
    Args:
        excel_bytes: Excel文档的字节数据
        
    Returns:
        PDF文档的字节数据
    """
    if pd is None or openpyxl is None:
        raise ImportError("pandas 或 openpyxl 未安装，无法处理Excel文档")
    
    try:
        # 使用BytesIO读取Excel文档，避免临时文件问题
        excel_stream = io.BytesIO(excel_bytes)
        
        # 读取Excel文档的所有工作表
        excel_file = pd.ExcelFile(excel_stream)
        
        # 创建PDF输出缓冲区
        pdf_buffer = io.BytesIO()
        
        # 创建PDF文档
        pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # 处理每个工作表
        for sheet_name in excel_file.sheet_names:
            # 添加工作表标题
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=16,
                alignment=TA_CENTER
            )
            story.append(Paragraph(f"工作表: {sheet_name}", title_style))
            story.append(Spacer(1, 0.2*inch))
            
            # 读取工作表数据
            df = pd.read_excel(excel_stream, sheet_name=sheet_name)
            
            if not df.empty:
                # 转换DataFrame为表格数据
                table_data = [df.columns.tolist()]  # 添加列标题
                for index, row in df.iterrows():
                    table_data.append([str(val) if pd.notna(val) else '' for val in row])
                
                # 创建PDF表格
                if len(table_data) > 0:
                    pdf_table = Table(table_data)
                    pdf_table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), 10),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black),
                        ('FONTSIZE', (0, 1), (-1, -1), 8),  # 数据行使用较小字体
                    ]))
                    story.append(pdf_table)
            else:
                story.append(Paragraph("工作表为空", styles['Normal']))
            
            story.append(Spacer(1, 0.3*inch))
        
        # 处理Excel文档中的图片（如果有的话）
        try:
            # 使用openpyxl重新加载工作簿以访问图片
            excel_stream.seek(0)  # 重置流位置
            from openpyxl import load_workbook
            wb = load_workbook(excel_stream)
            
            excel_images = []
            for sheet in wb.worksheets:
                if hasattr(sheet, '_images'):
                    for img in sheet._images:
                        if hasattr(img, '_data'):
                            excel_images.append(img._data())
            
            if excel_images and Image:
                # 添加图片部分标题
                story.append(Paragraph("Excel图片", styles['Heading2']))
                story.append(Spacer(1, 0.1*inch))
                
                for i, image_data in enumerate(excel_images):
                    try:
                        image_stream = io.BytesIO(image_data)
                        pil_image = Image.open(image_stream)
                        
                        # 保持较高的图片分辨率以确保OCR质量
                        max_width = 1200
                        max_height = 900
                        
                        # 只在图片过大时才缩放
                        if pil_image.width > max_width or pil_image.height > max_height:
                            pil_image.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)
                        
                        # 增强图片质量以提高OCR效果
                        pil_image = _enhance_image_for_ocr(pil_image)
                        
                        # 保存增强后的图片
                        resized_stream = io.BytesIO()
                        pil_image.save(resized_stream, format='PNG', quality=95, optimize=True)
                        resized_stream.seek(0)
                        
                        # 计算适合PDF页面的图片尺寸
                        page_width = 6*inch  # A4页面宽度减去边距
                        page_height = 8*inch  # A4页面高度减去边距
                        
                        # 计算缩放比例以适应页面
                        width_ratio = page_width / pil_image.width
                        height_ratio = page_height / pil_image.height
                        scale_ratio = min(width_ratio, height_ratio, 1.0)  # 不放大图片
                        
                        final_width = pil_image.width * scale_ratio
                        final_height = pil_image.height * scale_ratio
                        
                        # 添加到PDF
                        rl_image = RLImage(resized_stream, width=final_width, height=final_height)
                        story.append(rl_image)
                        story.append(Spacer(1, 0.2*inch))
                    except Exception as e:
                        logger.warning(f"处理Excel图片 {i+1} 失败: {e}")
                        
        except Exception as e:
            logger.debug(f"Excel图片提取失败: {e}")
        
        # 如果没有内容，添加一个提示
        if not story:
            story.append(Paragraph("Excel文档内容为空或无法解析", styles['Normal']))
        
        # 生成PDF
        pdf_doc.build(story)
        
        return pdf_buffer.getvalue()
            
    except Exception as e:
        logger.error(f"Excel转PDF转换失败: {e}")
        raise


def ppt_to_pdf_bytes(ppt_bytes: bytes) -> bytes:
    """
    将PowerPoint文档字节转换为PDF字节
    
    Args:
        ppt_bytes: PowerPoint文档的字节数据
        
    Returns:
        PDF文档的字节数据
    """
    if Presentation is None:
        raise ImportError("python-pptx 未安装，无法处理PowerPoint文档")
    
    try:
        # 使用BytesIO读取PowerPoint文档，避免临时文件问题
        ppt_stream = io.BytesIO(ppt_bytes)
        
        # 读取PowerPoint文档
        prs = Presentation(ppt_stream)
        
        # 创建PDF输出缓冲区
        pdf_buffer = io.BytesIO()
        
        # 创建PDF文档
        pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # 处理每张幻灯片
        for slide_idx, slide in enumerate(prs.slides, 1):
            # 添加幻灯片标题
            title_style = ParagraphStyle(
                'SlideTitle',
                parent=styles['Heading1'],
                fontSize=14,
                alignment=TA_CENTER
            )
            story.append(Paragraph(f"幻灯片 {slide_idx}", title_style))
            story.append(Spacer(1, 0.1*inch))
            
            # 提取幻灯片文本内容
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            # 添加文本内容到PDF
            if slide_text:
                for text in slide_text:
                    para = Paragraph(text, styles['Normal'])
                    story.append(para)
                    story.append(Spacer(1, 0.1*inch))
            else:
                story.append(Paragraph("此幻灯片无文本内容", styles['Normal']))
            
            # 添加分页符（除了最后一张幻灯片）
            if slide_idx < len(prs.slides):
                story.append(Spacer(1, 0.3*inch))
                story.append(Paragraph("─" * 50, styles['Normal']))
                story.append(Spacer(1, 0.3*inch))
        
        # 如果没有内容，添加一个提示
        if not story:
            story.append(Paragraph("PowerPoint文档内容为空或无法解析", styles['Normal']))
        
        # 生成PDF
        pdf_doc.build(story)
        
        return pdf_buffer.getvalue()
            
    except Exception as e:
        logger.error(f"PowerPoint转PDF转换失败: {e}")
        raise


def office_to_pdf_bytes(office_bytes: bytes, file_extension: str) -> bytes:
    """
    将Office文档字节转换为PDF字节
    
    Args:
        office_bytes: Office文档的字节数据
        file_extension: 文件扩展名 (.docx, .xlsx, .pptx等)
        
    Returns:
        PDF文档的字节数据
        
    Raises:
        ValueError: 不支持的文件格式
    """
    file_extension = file_extension.lower()
    
    if file_extension == '.docx':
        # 对于 .docx，优先直接内容提取并为图片生成高质量PDF
        logger.info("检测到Word文档（.docx），建议使用直接处理方式而非PDF转换")
        return word_to_pdf_bytes(office_bytes)
    elif file_extension == '.doc':
        # 对于旧版 .doc（非ZIP容器），python-docx 无法解析，需走系统级转换
        logger.info("检测到旧版Word文档（.doc），尝试系统级转换为PDF")
        pdf_bytes = _convert_legacy_doc_to_pdf_bytes(office_bytes)
        return pdf_bytes
    elif file_extension in ['.xlsx', '.xls']:
        # 对于 Excel，建议使用直接处理方式而非PDF转换
        logger.info("检测到Excel文档，建议使用直接处理方式而非PDF转换")
        return excel_to_pdf_bytes(office_bytes)
    elif file_extension in ['.pptx', '.ppt']:
        # 对于 PowerPoint，建议使用直接处理方式而非PDF转换
        logger.info("检测到PowerPoint文档，建议使用直接处理方式而非PDF转换")
        return ppt_to_pdf_bytes(office_bytes)
    else:
        raise ValueError(f"不支持的Office文档格式: {file_extension}")


def _convert_legacy_doc_to_pdf_bytes(doc_bytes: bytes) -> bytes:
    """
    将旧版 .doc 文档转换为 PDF 字节：
    - Windows: 优先使用 win32com（需安装 Microsoft Word 与 pywin32）
    - 其它平台或无 Word: 尝试调用 libreoffice/soffice CLI（需本机安装）
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        doc_path = os.path.join(tmpdir, 'input.doc')
        pdf_path = os.path.join(tmpdir, 'output.pdf')
        with open(doc_path, 'wb') as f:
            f.write(doc_bytes)

        # Windows 下优先使用 Word COM 自动化
        if platform_system().lower() == 'windows':
            try:
                import win32com.client  # type: ignore
                word = win32com.client.Dispatch('Word.Application')
                word.Visible = False
                try:
                    doc = word.Documents.Open(doc_path)
                    # 17: wdFormatPDF
                    doc.SaveAs(pdf_path, FileFormat=17)
                    doc.Close(False)
                finally:
                    word.Quit()
                with open(pdf_path, 'rb') as f:
                    return f.read()
            except Exception as e:
                logger.warning(f"win32com 转换 .doc 失败，尝试使用 libreoffice：{e}")

        # 尝试使用 libreoffice/soffice
        soffice = shutil.which('soffice') or shutil.which('soffice.exe')
        if not soffice:
            raise RuntimeError(
                "无法转换 .doc：未检测到 Microsoft Word（win32com）或 libreoffice(soffice)。"
            )
        try:
            # --headless --convert-to pdf --outdir <dir> <file>
            cmd = [
                soffice,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", tmpdir,
                doc_path,
            ]
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            # libreoffice 默认输出同名 pdf
            guessed_pdf = os.path.splitext(doc_path)[0] + '.pdf'
            if os.path.exists(guessed_pdf):
                with open(guessed_pdf, 'rb') as f:
                    return f.read()
            elif os.path.exists(pdf_path):
                with open(pdf_path, 'rb') as f:
                    return f.read()
            else:
                raise RuntimeError("libreoffice 转换未生成 PDF 文件")
        except Exception as e:
            raise RuntimeError(f"libreoffice 转换 .doc 失败: {e}")


def process_word_document_content(word_bytes: bytes) -> str:
    """
    直接提取Word文档的文本内容，返回Markdown格式
    这是一个新的处理方式，不经过PDF转换
    
    Args:
        word_bytes: Word文档的字节数据
        
    Returns:
        str: Markdown格式的文档内容
    """
    try:
        from doculook.utils.word_direct_processor import process_word_document_directly
        
        result = process_word_document_directly(word_bytes)
        
        # 合并文本内容和图片信息
        markdown_content = result["markdown_content"]
        
        if result["has_images"]:
            # 图片需要单独处理OCR，这里先添加占位符
            logger.info(f"Word文档包含 {result['image_count']} 张图片，需要进行OCR处理")
            # 实际的图片OCR结果需要在后续流程中替换这些占位符
        
        return markdown_content
        
    except Exception as e:
        logger.error(f"Word文档直接内容提取失败: {e}")
        # 回退到原有的PDF转换方式
        logger.info("回退到PDF转换方式处理Word文档")
        raise


def process_excel_document_content(excel_bytes: bytes) -> str:
    """
    直接提取Excel文档的内容，返回Markdown格式
    这是一个新的处理方式，不经过PDF转换
    
    Args:
        excel_bytes: Excel文档的字节数据
        
    Returns:
        str: Markdown格式的文档内容
    """
    try:
        from doculook.utils.excel_direct_processor import process_excel_document_directly
        
        result = process_excel_document_directly(excel_bytes)
        
        # 获取Markdown内容
        markdown_content = result["markdown_content"]
        
        if result["has_images"]:
            # 图片需要单独处理OCR，这里先添加占位符
            logger.info(f"Excel文档包含 {result['image_count']} 张图片，需要进行OCR处理")
            # 实际的图片OCR结果需要在后续流程中替换这些占位符
        
        return markdown_content
        
    except Exception as e:
        logger.error(f"Excel文档直接内容提取失败: {e}")
        # 回退到原有的PDF转换方式
        logger.info("回退到PDF转换方式处理Excel文档")
        raise


def process_ppt_document_content(ppt_bytes: bytes) -> str:
    """
    直接提取PowerPoint文档的内容，返回Markdown格式
    这是一个新的处理方式，不经过PDF转换
    
    Args:
        ppt_bytes: PowerPoint文档的字节数据
        
    Returns:
        str: Markdown格式的文档内容
    """
    try:
        from doculook.utils.ppt_direct_processor import process_ppt_document_directly
        
        result = process_ppt_document_directly(ppt_bytes)
        
        # 获取Markdown内容
        markdown_content = result["markdown_content"]
        
        if result["has_images"]:
            # 图片需要单独处理OCR，这里先添加占位符
            logger.info(f"PowerPoint文档包含 {result['image_count']} 张图片，需要进行OCR处理")
            # 实际的图片OCR结果需要在后续流程中替换这些占位符
        
        return markdown_content
        
    except Exception as e:
        logger.error(f"PowerPoint文档直接内容提取失败: {e}")
        # 回退到原有的PDF转换方式
        logger.info("回退到PDF转换方式处理PowerPoint文档")
        raise


# 支持的Office文档后缀
OFFICE_SUFFIXES = ['.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt']


def is_office_file(file_path: Union[str, Path]) -> bool:
    """
    检查文件是否为支持的Office文档格式
    
    Args:
        file_path: 文件路径
        
    Returns:
        bool: 如果是支持的Office文档格式返回True，否则返回False
    """
    if isinstance(file_path, str):
        file_path = Path(file_path)
    
    return file_path.suffix.lower() in OFFICE_SUFFIXES


if __name__ == "__main__":
    # 测试代码
    logger.info("Office转换器模块加载完成")
    logger.info(f"支持的Office文档格式: {OFFICE_SUFFIXES}")
