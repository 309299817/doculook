#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint文档直接处理器
直接从PowerPoint文档中提取内容，包括文本、图片、表格等，无需转换为PDF
"""

import io
import json
from typing import Dict, List, Tuple, Any, Optional
from pathlib import Path
from loguru import logger

try:
    from pptx import Presentation
    from pptx.shapes.picture import Picture
    from pptx.shapes.group import GroupShape
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    # 注意：Table 不是一个独立的导入类，而是通过 shape.table 属性访问
    Table = None  # 占位符，实际使用时通过 shape.table 访问
except ImportError:
    logger.warning("python-pptx 未安装，PowerPoint直接处理功能将被禁用")
    Presentation = None
    Picture = None
    Table = None
    GroupShape = None
    MSO_SHAPE_TYPE = None

try:
    from PIL import Image
except ImportError:
    logger.warning("PIL(Pillow) 未安装，图片处理功能受限")
    Image = None


def _clean_text_content(text: str) -> str:
    """
    清理文本内容，处理编码和显示问题
    
    Args:
        text: 原始文本
        
    Returns:
        str: 清理后的文本
    """
    if not text:
        return ""
    
    try:
        # 1. 处理编码问题
        if isinstance(text, bytes):
            text = text.decode('utf-8', errors='replace')
        
        # 2. 替换常见的问题字符
        # 方块字符通常是字体缺失或编码问题
        problem_chars = {
            '■': '',  # 移除方块字符
            '□': '',  # 移除空心方块
            '\ufffd': '',  # 移除替换字符
            '\u25a0': '',  # 移除黑色方块
            '\u25a1': '',  # 移除白色方块
        }
        
        for problem_char, replacement in problem_chars.items():
            text = text.replace(problem_char, replacement)
        
        # 3. 清理控制字符，但保留常用的空白字符
        cleaned_text = ''.join(
            char for char in text 
            if ord(char) >= 32 or char in '\n\r\t'
        )
        
        # 4. 清理多余的空白
        import re
        cleaned_text = re.sub(r'\s+', ' ', cleaned_text).strip()
        
        # 5. 如果清理后为空，记录警告
        if not cleaned_text and text:
            logger.warning(f"文本清理后变为空: 原文本长度={len(text)}")
        
        return cleaned_text
        
    except Exception as e:
        logger.warning(f"文本清理失败: {e}")
        return text  # 返回原始文本


def _extract_text_from_shape(shape) -> str:
    """
    从形状中提取文本
    
    Args:
        shape: PowerPoint形状对象
        
    Returns:
        str: 提取的文本
    """
    try:
        # 方法1: 直接获取text属性
        if hasattr(shape, 'text') and shape.text and shape.text.strip():
            text_content = shape.text.strip()
            if text_content:
                # 处理可能的编码问题
                try:
                    # 确保文本是正确的Unicode字符串
                    if isinstance(text_content, bytes):
                        text_content = text_content.decode('utf-8', errors='replace')
                    
                    # 清理可能的控制字符
                    text_content = ''.join(char for char in text_content if ord(char) >= 32 or char in '\n\r\t')
                    
                    # 应用文本清理
                    text_content = _clean_text_content(text_content)
                    
                    logger.debug(f"通过shape.text提取到文本: {text_content[:50]}...")
                    return text_content
                except Exception as encoding_error:
                    logger.warning(f"文本编码处理失败: {encoding_error}")
                    return text_content  # 返回原始文本
        
        # 方法2: 通过text_frame深度提取
        if hasattr(shape, 'text_frame') and shape.text_frame:
            all_texts = []
            try:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_texts = []
                    for run in paragraph.runs:
                        if run.text and run.text.strip():
                            paragraph_texts.append(run.text.strip())
                    
                    if paragraph_texts:
                        paragraph_text = ' '.join(paragraph_texts)
                        all_texts.append(paragraph_text)
                
                if all_texts:
                    combined_text = '\n'.join(all_texts)
                    # 应用文本清理
                    combined_text = _clean_text_content(combined_text)
                    logger.debug(f"通过text_frame提取到文本: {combined_text[:50]}...")
                    return combined_text
                    
            except Exception as frame_error:
                logger.debug(f"text_frame处理失败: {frame_error}")
        
        # 方法3: 检查是否是占位符
        if hasattr(shape, 'placeholder_format'):
            try:
                placeholder = shape.placeholder_format
                if placeholder and hasattr(placeholder, 'element'):
                    # 这可能是标题或内容占位符
                    logger.debug(f"发现占位符形状，尝试提取文本")
                    if hasattr(shape, 'text') and shape.text:
                        return shape.text.strip()
            except Exception as placeholder_error:
                logger.debug(f"占位符处理失败: {placeholder_error}")
        
        # 方法4: 检查是否有table属性（表格中的文本）
        if hasattr(shape, 'table'):
            try:
                table_texts = []
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        cell_text = _extract_text_from_shape(cell)
                        if cell_text.strip():
                            table_texts.append(cell_text.strip())
                
                if table_texts:
                    logger.debug(f"从表格中提取到文本: {len(table_texts)}个单元格")
                    return ' | '.join(table_texts)  # 简单的表格文本连接
                    
            except Exception as table_error:
                logger.debug(f"表格文本提取失败: {table_error}")
        
    except Exception as e:
        logger.debug(f"提取形状文本失败: {e}")
    
    return ""


def _extract_table_from_shape(shape) -> str:
    """
    从表格形状中提取表格数据并转换为Markdown
    
    Args:
        shape: PowerPoint表格形状对象
        
    Returns:
        str: Markdown格式的表格
    """
    if not hasattr(shape, 'table'):
        return ""
    
    try:
        table = shape.table
        table_data = []
        
        # 提取表格数据
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = _extract_text_from_shape(cell)
                # 应用文本清理
                cell_text = _clean_text_content(cell_text)
                # 转义Markdown特殊字符
                cell_text = cell_text.replace('|', '&#124;').replace('\n', '<br/>')
                row_data.append(cell_text or '')
            table_data.append(row_data)
        
        if not table_data:
            return ""
        
        # 生成Markdown表格
        markdown_lines = []
        
        # 表头
        if len(table_data) > 0:
            header_row = "| " + " | ".join(table_data[0]) + " |"
            markdown_lines.append(header_row)
            
            # 分隔符行
            separator_row = "| " + " | ".join(["---"] * len(table_data[0])) + " |"
            markdown_lines.append(separator_row)
            
            # 数据行
            for row in table_data[1:]:
                # 确保行的列数与表头一致
                while len(row) < len(table_data[0]):
                    row.append('')
                data_row = "| " + " | ".join(row[:len(table_data[0])]) + " |"
                markdown_lines.append(data_row)
        
        return '\n'.join(markdown_lines)
        
    except Exception as e:
        logger.warning(f"提取表格失败: {e}")
        return ""


def _extract_image_from_shape(shape, image_index: int) -> Optional[Dict]:
    """
    从图片形状中提取图片数据
    
    Args:
        shape: PowerPoint图片形状对象
        image_index: 图片索引
        
    Returns:
        Optional[Dict]: 图片信息字典
    """
    try:
        if hasattr(shape, 'image') and shape.image:
            image_part = shape.image
            image_bytes = image_part.blob
            
            # 验证图片数据有效性
            if not image_bytes or len(image_bytes) < 100:
                logger.warning(f"图片 {image_index} 数据过小或为空")
                return None
            
            # 尝试检测图片格式
            image_format = "png"  # 默认格式
            if image_bytes.startswith(b'\xff\xd8\xff'):
                image_format = "jpg"
            elif image_bytes.startswith(b'\x89PNG'):
                image_format = "png"
            elif image_bytes.startswith(b'GIF8'):
                image_format = "gif"
            elif image_bytes.startswith(b'BM'):
                image_format = "bmp"
            
            # 使用PIL验证图片完整性
            if Image:
                try:
                    import io
                    img_stream = io.BytesIO(image_bytes)
                    test_img = Image.open(img_stream)
                    test_img.verify()  # 验证图片完整性
                    logger.debug(f"图片 {image_index} 格式验证成功: {image_format}, 大小: {len(image_bytes)} bytes")
                except Exception as verify_error:
                    logger.warning(f"图片 {image_index} 验证失败: {verify_error}")
                    return None
            
            return {
                'data': image_bytes,
                'format': image_format,
                'width': getattr(shape, 'width', None),
                'height': getattr(shape, 'height', None),
                'left': getattr(shape, 'left', None),
                'top': getattr(shape, 'top', None),
                'size': len(image_bytes)
            }
            
    except Exception as e:
        logger.warning(f"提取图片 {image_index} 失败: {e}")
    return None


def _process_shape_group(shape_group, slide_index: int, all_images: List, image_counter: List[int]) -> str:
    """
    递归处理形状组
    
    Args:
        shape_group: 形状组对象
        slide_index: 幻灯片索引
        all_images: 所有图片数据列表
        image_counter: 图片计数器（使用列表实现引用传递）
        
    Returns:
        str: 处理结果的Markdown内容
    """
    content_parts = []
    
    for shape in shape_group.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # 递归处理嵌套的形状组
            group_content = _process_shape_group(shape, slide_index, all_images, image_counter)
            if group_content.strip():
                content_parts.append(group_content)
        else:
            # 处理普通形状
            shape_content = _process_single_shape(shape, slide_index, all_images, image_counter)
            if shape_content.strip():
                content_parts.append(shape_content)
    
    return '\n\n'.join(content_parts)


def _process_single_shape(shape, slide_index: int, all_images: List, image_counter: List[int]) -> str:
    """
    处理单个形状
    
    Args:
        shape: 形状对象
        slide_index: 幻灯片索引
        all_images: 所有图片数据列表
        image_counter: 图片计数器
        
    Returns:
        str: 处理结果的Markdown内容
    """
    content_parts = []
    
    try:
        # 记录形状信息用于调试
        shape_type = getattr(shape, 'shape_type', None)
        shape_name = getattr(shape, 'name', 'Unknown')
        logger.debug(f"处理形状: {shape_name}, 类型: {shape_type}")
        
        # 1. 优先处理图片
        if MSO_SHAPE_TYPE and shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_counter[0] += 1
            image_info = _extract_image_from_shape(shape, image_counter[0])
            if image_info:
                all_images.append(image_info['data'])
                # 使用检测到的格式作为文件扩展名
                image_format = image_info.get('format', 'png')
                image_ref = f"![图片{image_counter[0]}](images/ppt_image_{image_counter[0]}.{image_format})"
                content_parts.append(image_ref)
                logger.info(f"幻灯片{slide_index+1}中发现图片 {image_counter[0]}: {image_info['size']} bytes, 格式: {image_format}")
        
        # 2. 处理表格
        elif MSO_SHAPE_TYPE and shape_type == MSO_SHAPE_TYPE.TABLE:
            table_md = _extract_table_from_shape(shape)
            if table_md.strip():
                content_parts.append(f"**表格:**\n{table_md}")
                logger.info(f"幻灯片{slide_index+1}中发现表格")
        
        # 3. 处理所有可能包含文本的形状
        else:
            # 尝试提取文本内容
            text_content = _extract_text_from_shape(shape)
            if text_content.strip():
                # 判断文本的重要性和格式
                is_title = False
                is_important = False
                
                # 检查是否是标题类型
                try:
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                        placeholder_type = shape.placeholder_format.type
                        if placeholder_type in [1, 3]:  # 标题或副标题
                            is_title = True
                        elif placeholder_type in [2, 7]:  # 内容或其他重要内容
                            is_important = True
                        logger.debug(f"占位符类型: {placeholder_type}")
                except:
                    pass
                
                # 检查文本框的位置和大小（标题通常在上方且较大）
                try:
                    if hasattr(shape, 'top') and hasattr(shape, 'height'):
                        # 如果形状在幻灯片上方1/3且高度较大，可能是标题
                        if shape.top < 2000000 and shape.height > 500000:  # PowerPoint使用EMU单位
                            is_title = True
                except:
                    pass
                
                # 根据判断结果格式化文本
                if is_title:
                    content_parts.append(f"### {text_content}")
                    logger.debug(f"识别为标题: {text_content[:30]}...")
                elif is_important:
                    content_parts.append(f"**{text_content}**")
                    logger.debug(f"识别为重要内容: {text_content[:30]}...")
                else:
                    content_parts.append(text_content)
                    logger.debug(f"识别为普通文本: {text_content[:30]}...")
        
        # 4. 特殊形状处理
        if not content_parts:
            # 如果没有提取到任何内容，记录形状信息
            logger.debug(f"幻灯片{slide_index+1}中的形状 {shape_name}(类型:{shape_type}) 未提取到内容")
    
    except Exception as e:
        logger.warning(f"处理幻灯片{slide_index+1}中的形状失败: {e}")
    
    return '\n\n'.join(content_parts)


def process_ppt_document_directly(ppt_bytes: bytes) -> Dict:
    """
    直接处理PowerPoint文档，提取文本、图片和表格
    
    Args:
        ppt_bytes: PowerPoint文档的字节数据
        
    Returns:
        Dict: 包含处理结果的字典
    """
    if Presentation is None:
        raise ImportError("python-pptx 未安装，无法处理PowerPoint文档")
    
    try:
        # 读取PowerPoint文档
        ppt_stream = io.BytesIO(ppt_bytes)
        prs = Presentation(ppt_stream)
        
        # 初始化结果
        result = {
            'markdown_content': '',
            'has_images': False,
            'image_count': 0,
            'has_tables': False,
            'has_data': False,
            'slides': []
        }
        
        markdown_parts = []
        all_images = []
        image_counter = [0]  # 使用列表实现引用传递
        
        # 添加文档标题
        markdown_parts.append("# PowerPoint文档内容\n")
        
        # 处理每张幻灯片
        for slide_index, slide in enumerate(prs.slides):
            slide_info = {
                'slide_number': slide_index + 1,
                'text_count': 0,
                'image_count': 0,
                'table_count': 0,
                'shape_count': len(slide.shapes)
            }
            
            slide_content = []
            slide_content.append(f"## 幻灯片 {slide_index + 1}")
            
            slide_has_content = False
            logger.info(f"开始处理幻灯片 {slide_index + 1}, 包含 {len(slide.shapes)} 个形状")
            
            # 处理幻灯片中的所有形状
            for shape in slide.shapes:
                shape_content = ""
                
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # 处理形状组
                    shape_content = _process_shape_group(shape, slide_index, all_images, image_counter)
                else:
                    # 处理单个形状
                    shape_content = _process_single_shape(shape, slide_index, all_images, image_counter)
                
                if shape_content.strip():
                    slide_content.append(shape_content)
                    slide_has_content = True
                    
                    # 统计内容类型
                    if '![图片' in shape_content:
                        slide_info['image_count'] += shape_content.count('![图片')
                    if '**表格:**' in shape_content:
                        slide_info['table_count'] += 1
                    if shape_content.strip() and not shape_content.startswith('![图片') and '**表格:**' not in shape_content:
                        slide_info['text_count'] += 1
            
            # 如果幻灯片没有内容，添加提示
            if not slide_has_content:
                slide_content.append("*此幻灯片无可提取的文本或媒体内容*")
            
            markdown_parts.extend(slide_content)
            markdown_parts.append("")  # 添加空行分隔
            
            result['slides'].append(slide_info)
        
        # 汇总结果
        result['markdown_content'] = '\n'.join(markdown_parts)
        result['image_count'] = len(all_images)
        result['has_images'] = len(all_images) > 0
        result['has_tables'] = any(slide['table_count'] > 0 for slide in result['slides'])
        result['has_data'] = True  # PowerPoint总是有数据（至少有幻灯片结构）
        
        # 保存图片数据和元数据
        result['images'] = all_images
        
        # 从all_images中提取元数据（注意：这里需要从原始图片信息中获取）
        # 由于当前实现只保存了图片字节数据，我们需要重新构建元数据
        image_metadata = []
        for i, image_data in enumerate(all_images):
            metadata = {'format': 'png'}  # 默认格式
            # 尝试检测格式
            if image_data.startswith(b'\xff\xd8\xff'):
                metadata['format'] = 'jpg'
            elif image_data.startswith(b'\x89PNG'):
                metadata['format'] = 'png'
            elif image_data.startswith(b'GIF8'):
                metadata['format'] = 'gif'
            elif image_data.startswith(b'BM'):
                metadata['format'] = 'bmp'
            
            metadata['size'] = len(image_data)
            image_metadata.append(metadata)
        
        result['image_metadata'] = image_metadata
        
        logger.info(f"PowerPoint文档处理完成: {len(prs.slides)}张幻灯片, {result['image_count']}张图片, {sum(slide['table_count'] for slide in result['slides'])}个表格")
        
        return result
        
    except Exception as e:
        logger.error(f"PowerPoint文档处理失败: {e}")
        raise


def is_ppt_document(file_path: Path) -> bool:
    """
    检查文件是否为PowerPoint文档
    
    Args:
        file_path: 文件路径
        
    Returns:
        bool: 是否为PowerPoint文档
    """
    return file_path.suffix.lower() in ['.pptx', '.ppt']
